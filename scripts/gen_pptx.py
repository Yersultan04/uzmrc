"""Генерация PPTX презентации UzMRC (MVP, рус)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "notes", "demo-shots")
OUT = os.path.join(HERE, "deliverables", "UzMRC-Презентация-MVP.pptx")
os.makedirs(os.path.dirname(OUT), exist_ok=True)

ACCENT = RGBColor(0x5B, 0x3D, 0xF5)
DARK = RGBColor(0x14, 0x14, 0x24)
LIGHT = RGBColor(0xF4, 0xF2, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x55, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, text, size=18, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def bullets(s, x, y, w, h, items, size=18, color=DARK, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "•  " + it
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def header(s, title, kicker=None):
    rect(s, 0, 0, SW, Inches(0.18), ACCENT)
    txt(s, Inches(0.7), Inches(0.5), Inches(12), Inches(1.0), title,
        size=32, color=DARK, bold=True)
    if kicker:
        txt(s, Inches(0.72), Inches(1.35), Inches(12), Inches(0.5), kicker,
            size=15, color=ACCENT, bold=True)


def kpi_card(s, x, y, w, h, value, label):
    rect(s, x, y, w, h, LIGHT)
    txt(s, x, y + Inches(0.25), w, Inches(0.9), value, size=34, color=ACCENT,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, y + Inches(1.15), w, Inches(0.5), label, size=13, color=GREY,
        align=PP_ALIGN.CENTER)


def pic_fit(s, path, x, y, max_w, max_h):
    if not os.path.exists(path):
        return None
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = (1200, 900)
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = int(max_w / ar)
    else:
        h = max_h; w = int(max_h * ar)
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    return s.shapes.add_picture(path, px, py, width=w, height=h)


# ---------- 1. Титул ----------
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(3.0), SW, Inches(0.06), ACCENT)
txt(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2), "UzMRC", size=66,
    color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.2), Inches(11.3), Inches(0.8),
    "ИИ-ассистент по нормативным документам", size=26, color=RGBColor(0xC9, 0xC2, 0xF5),
    align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
    "Чат по документам со ссылкой на источник  ·  сравнение приказов с нормами",
    size=16, color=RGBColor(0x9A, 0x96, 0xB5), align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
    "MVP для демонстрации  ·  18 июня 2026", size=13,
    color=RGBColor(0x77, 0x74, 0x90), align=PP_ALIGN.CENTER)

# ---------- 2. Проблема ----------
s = slide()
header(s, "Задача", "Зачем это нужно")
bullets(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(4.5), [
    "Нормативная база компании — десятки документов на русском и узбекском. Найти нужную норму и точную формулировку вручную — долго.",
    "Новые приказы и регламенты нужно сверять с действующими нормами: не противоречит ли, не дублирует ли, нет ли пробела.",
    "Цена ошибки высока — нужна ссылка на источник и точная цитата, а не «пересказ» от ИИ.",
    "Решение: ассистент, который отвечает строго по базе со ссылкой на файл/страницу/цитату и сам сравнивает документы с нормами.",
], size=19, gap=14)

# ---------- 3. Что умеет ----------
s = slide()
header(s, "Что умеет система", "Два модуля")
rect(s, Inches(0.7), Inches(2.1), Inches(5.85), Inches(4.4), LIGHT)
txt(s, Inches(1.0), Inches(2.35), Inches(5.3), Inches(0.6), "1. Чат по документам",
    size=22, color=ACCENT, bold=True)
bullets(s, Inches(1.0), Inches(3.1), Inches(5.3), Inches(3.2), [
    "Вопрос на рус/узб → ответ по базе",
    "Ссылка: файл, страница, точная цитата",
    "Кросс-языковой поиск ru ↔ uz",
    "Нет ответа в базе → не выдумывает",
], size=16, gap=10)
rect(s, Inches(6.8), Inches(2.1), Inches(5.85), Inches(4.4), LIGHT)
txt(s, Inches(7.1), Inches(2.35), Inches(5.3), Inches(0.6), "2. Сравнение документов",
    size=22, color=ACCENT, bold=True)
bullets(s, Inches(7.1), Inches(3.1), Inches(5.3), Inches(3.2), [
    "Приказ разбивается на пункты",
    "Каждый пункт сверяется с базой",
    "Противоречие / дубль / пробел / дополнение",
    "Обоснование + цитата нормы + рекомендация",
    "Экспорт отчёта (.md / PDF)",
], size=16, gap=10)

# ---------- 4. Модуль 1 ----------
s = slide()
header(s, "Модуль 1 — Чат по документам", "Агентный RAG")
bullets(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(4.5), [
    "Гибридный поиск: плотные векторы (Voyage + Qdrant) + полнотекстовый (Postgres), слияние RRF, LLM-rerank.",
    "Агент ReAct: цикл рассуждение→поиск→наблюдение, до 40 шагов, набор инструментов (поиск, декомпозиция, точный lookup, выборка страниц).",
    "Структурированный вывод: каждый шаг — валидируемый JSON (защита от «свободного» текста модели).",
    "Grounding-проверка: каждая цитата сверяется с исходным чанком; уверенность ограничивается долей подтверждённых цитат.",
], size=18, gap=14)

# ---------- 5. Модуль 2 ----------
s = slide()
header(s, "Модуль 2 — Сравнение документов", "Типы находок")
data = [
    ("Противоречие", "пункт нельзя исполнить, не нарушив действующую норму", RGBColor(0xE2,0x4A,0x4A)),
    ("Дубль", "повторяет уже существующую норму", RGBColor(0x3D,0x8B,0xF5)),
    ("Дополнение", "расширяет/уточняет, не противореча", RGBColor(0x2E,0xA0,0x6A)),
    ("Пробел", "в базе нет нормы по теме", RGBColor(0xE0,0x8A,0x1E)),
]
y = Inches(2.2)
for name, desc, col in data:
    rect(s, Inches(0.7), y, Inches(0.22), Inches(0.95), col)
    txt(s, Inches(1.1), y, Inches(3.0), Inches(0.95), name, size=20, color=DARK,
        bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.3), y, Inches(8.2), Inches(0.95), desc, size=17, color=GREY,
        anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.05)
txt(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.6),
    "Фоновый прогон с прогресс-баром · ~20–25 с на 10 пунктов · до 120 пунктов за прогон",
    size=14, color=ACCENT, bold=True)

# ---------- 6. Архитектура ----------
s = slide()
header(s, "Архитектура и стек", "Docker Compose · 4 контейнера")
rows = [
    ("Frontend", "React + Vite + TypeScript + Tailwind, nginx"),
    ("Backend", "FastAPI (Python 3.12), async SQLAlchemy, Pydantic v2"),
    ("Реляционная БД", "PostgreSQL 16 — метаданные, чанки, FTS, кэш эмбеддингов"),
    ("Векторная БД", "Qdrant — плотные векторы (cosine)"),
    ("Эмбеддинги", "Voyage voyage-3.5 (1024-dim, multilingual)"),
    ("LLM-судья / rerank", "Cerebras gpt-oss-120b, фолбэк OpenRouter"),
    ("Авторизация", "JWT (HS256), bcrypt, закрытая регистрация"),
]
y = Inches(2.05)
for i, (k, v) in enumerate(rows):
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.7), y, Inches(11.9), Inches(0.62), bg)
    txt(s, Inches(0.9), y, Inches(3.4), Inches(0.62), k, size=15, color=ACCENT,
        bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.4), y, Inches(8.0), Inches(0.62), v, size=15, color=DARK,
        anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.64)

# ---------- 7. Анти-галлюцинации ----------
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, SW, Inches(0.18), ACCENT)
txt(s, Inches(0.7), Inches(0.6), Inches(12), Inches(1), "Защита от галлюцинаций",
    size=34, color=WHITE, bold=True)
txt(s, Inches(0.72), Inches(1.5), Inches(12), Inches(0.5),
    "Ключевое требование для нормативного ассистента", size=16,
    color=RGBColor(0xC9,0xC2,0xF5), bold=True)
bullets(s, Inches(0.7), Inches(2.4), Inches(11.9), Inches(3.5), [
    "Каждая цитата сверяется с исходным фрагментом: подстрока → нормализация (переносы, пунктуация, гомоглифы OCR) → fuzzy-сопоставление.",
    "Уверенность ответа ограничивается долей подтверждённых цитат: не подтвердились 2 из 3 — уверенность не выше 0.33.",
    "Нет нормы в базе → система честно не отвечает, а не фабрикует ссылку (проверено на вопросе вне базы).",
], size=19, color=WHITE, gap=16)
rect(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.9), ACCENT)
txt(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.9),
    "Лучше честно не ответить, чем сослаться на несуществующую норму",
    size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------- 8. База знаний (цифры + скрин) ----------
s = slide()
header(s, "База знаний — текущий корпус", "Живые цифры на странице «О системе»")
kw = Inches(2.75); kh = Inches(1.75); ky = Inches(2.1); kx = Inches(0.7); gapx = Inches(0.18)
kpis = [("50", "документов"), ("1 630", "чанков"), ("32.6", "чанков/док"), ("643 974", "токенов")]
for i, (v, l) in enumerate(kpis):
    kpi_card(s, kx + i * (kw + gapx), ky, kw, kh, v, l)
txt(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(0.6),
    "Источник: ≈35 нормативных актов + ≈12 аналитических обзоров (uzmrc.uz, ≈595 стр.) · voyage-3.5 · статус «Готова»",
    size=14, color=GREY)
pic_fit(s, os.path.join(SHOTS, "uzmrc-about-page.png"), Inches(3.4), Inches(4.7),
        Inches(6.5), Inches(2.6))

# ---------- 9. Скрин отчёта сравнения ----------
s = slide()
header(s, "Пример отчёта сравнения", "Демо-приказ: 10 пунктов")
txt(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.5),
    "4 противоречия · 1 дубль · 1 дополнение · 4 пробела · цитаты норм подтверждены 5/7 · ~23 с",
    size=15, color=ACCENT, bold=True)
pic_fit(s, os.path.join(SHOTS, "uzmrc-compare-report-fullpage.png"), Inches(2.4),
        Inches(2.5), Inches(8.5), Inches(4.7))

# ---------- 10. MVP scope ----------
s = slide()
header(s, "Границы MVP", "Что входит и что нет")
rect(s, Inches(0.7), Inches(2.1), Inches(5.85), Inches(4.5), RGBColor(0xEA,0xF7,0xEF))
txt(s, Inches(1.0), Inches(2.3), Inches(5.3), Inches(0.5), "Входит", size=20,
    color=RGBColor(0x2E,0xA0,0x6A), bold=True)
bullets(s, Inches(1.0), Inches(3.0), Inches(5.3), Inches(3.4), [
    "Чат по базе со ссылками на источник",
    "Сравнение документа с нормами",
    "Кросс-язык ru ↔ uz",
    "Экспорт отчёта (.md / PDF)",
], size=16, gap=12)
rect(s, Inches(6.8), Inches(2.1), Inches(5.85), Inches(4.5), RGBColor(0xFB,0xEC,0xEC))
txt(s, Inches(7.1), Inches(2.3), Inches(5.3), Inches(0.5), "Не входит / ограничения",
    size=20, color=RGBColor(0xD0,0x3D,0x3D), bold=True)
bullets(s, Inches(7.1), Inches(3.0), Inches(5.3), Inches(3.4), [
    "Анализ портфеля, маскирование ПДн",
    "Генерация презентаций, парсинг рынка/OLX",
    "До 120 пунктов за прогон сравнения",
    "Загрузка новых док — интернет + платный Voyage",
    "Демо-набор 50 док, не полная база",
], size=16, gap=12)

# ---------- 11. Метрики ----------
s = slide()
header(s, "Качество", "Проверено вживую")
bullets(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(4.4), [
    "Кросс-язык: русский запрос «антикоррупционная политика» находит русский документ (score 0.80) и узбекские (0.79).",
    "Сравнение демо-приказа: 10 пунктов → 4 противоречия / 1 дубль / 1 дополнение / 4 пробела, цитаты норм 5/7, ~23 с.",
    "Скорость: чат по готовой базе — секунды; сравнение 10 пунктов — ~20–25 с; реиндекс с кэшем эмбеддингов — ×77 быстрее.",
    "Eval-наборы (retrieval / compare) + baseline зафиксированы в репозитории; 10 реальных примеров оформлены отдельно.",
], size=18, gap=16)

# ---------- 12. Развёртывание ----------
s = slide()
header(s, "Развёртывание", "Docker · локально или на сервере")
bullets(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(3.0), [
    "Локально: docker compose -f docker-compose.prod.yml up -d --build → UI на http://localhost:8090",
    "Постоянный публичный адрес: деплой на VPS (домен + reverse-proxy).",
    "Временный доступ: Cloudflare quick tunnel — бесплатно, без регистрации (URL эфемерный).",
    "Документация: QUICKSTART.md (запуск) и notes/DEPLOY-RUNBOOK.md (деплой).",
], size=18, gap=14)

# ---------- 13. Финал ----------
s = slide()
rect(s, 0, 0, SW, SH, DARK)
txt(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
    "UzMRC — готов к демонстрации", size=40, color=WHITE, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.8),
    "Оба модуля работают end-to-end · 50 документов в базе · ссылки на источник",
    size=18, color=RGBColor(0xC9,0xC2,0xF5), align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
    "github.com/Yersultan04/uzmrc", size=15, color=RGBColor(0x9A,0x96,0xB5),
    align=PP_ALIGN.CENTER)

prs.save(OUT)
print("OK ->", OUT, "| слайдов:", len(prs.slides._sldIdLst))
